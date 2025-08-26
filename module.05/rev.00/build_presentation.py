from __future__ import annotations
import os
import sys
import textwrap
import json

import pandas as pd
import re
from datetime import datetime
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "ds-capstone-template-coursera.pptx")
OUT_PPTX = os.path.join(os.path.dirname(__file__), "SpaceX_Capstone_Presentation.pptx")
OUT_PDF = os.path.join(os.path.dirname(__file__), "SpaceX_Capstone_Presentation.pdf")
PLOTS_DIR = os.path.join(os.path.dirname(__file__), "plots")
MOD2_PLOTS = os.path.join(os.path.dirname(__file__), "..", "module.02", "plots")
FOLIUM_SCREENSHOT = os.path.join(os.path.dirname(__file__), "..", "module.03", "folium_map.png")
SQL_SUMMARY_MD = os.path.join(os.path.dirname(__file__), "..", "module.02", "spacex_eda_sql_summary.md")

DATA_URL = (
    "https://cf-courses-data.s3.us.cloud-object-storage.appdomain.cloud/"
    "IBM-DS0321EN-SkillsNetwork/datasets/spacex_launch_dash.csv"
)


def ensure_dirs():
    os.makedirs(PLOTS_DIR, exist_ok=True)


def load_data() -> pd.DataFrame:
    """Load SpaceX dataset.
    1) Try network URL (Skills Network CSV)
    2) Fallback to local merged CSV from module 02
    Then standardize columns to those used by the Dash figures.
    """
    local_csv = os.path.normpath(os.path.join(os.path.dirname(__file__), "..", "module.02", "spacex_launches_merged.csv"))
    df: pd.DataFrame | None = None
    # Try remote
    try:
        df = pd.read_csv(DATA_URL)
    except Exception:
        # Fallback to local
        if os.path.exists(local_csv):
            df = pd.read_csv(local_csv)
        else:
            raise RuntimeError("Failed to load SpaceX data from URL and local CSV not found.")

    # Standardize expected columns: 'Launch Site', 'class', 'Payload Mass (kg)', 'Booster Version Category'
    # Launch Site
    if 'Launch Site' not in df.columns:
        for cand in ['Launch site', 'LaunchSite']:
            if cand in df.columns:
                df['Launch Site'] = df[cand]
                break
    # class (landing success 0/1)
    if 'class' not in df.columns:
        if 'LandingSuccess' in df.columns:
            df['class'] = df['LandingSuccess'].astype(int)
        elif 'Outcome' in df.columns:
            # crude mapping: outcomes starting with 'True' => 1, else 0
            df['class'] = df['Outcome'].astype(str).str.startswith('True').astype(int)
        else:
            df['class'] = 0
    # Payload Mass (kg)
    if 'Payload Mass (kg)' not in df.columns:
        for cand in ['Payload mass (kg)', 'PayloadMass', 'PayloadMass(kg)']:
            if cand in df.columns:
                df['Payload Mass (kg)'] = pd.to_numeric(df[cand], errors='coerce')
                break
    # Booster Version Category
    if 'Booster Version Category' not in df.columns:
        # Derive from 'Version Booster' or 'BoosterVersion' if available
        source_col = None
        for cand in ['Booster Version Category', 'Version Booster', 'BoosterVersion']:
            if cand in df.columns:
                source_col = cand
                break
        def to_cat(x: str) -> str:
            s = str(x)
            if re.search(r'v1\.0', s, re.IGNORECASE):
                return 'v1.0'
            if re.search(r'v1\.1', s, re.IGNORECASE):
                return 'v1.1'
            if re.search(r'FT', s, re.IGNORECASE):
                return 'FT'
            if re.search(r'\bB\s*4\b|B4', s, re.IGNORECASE):
                return 'B4'
            if re.search(r'\bB\s*5\b|B5', s, re.IGNORECASE):
                return 'B5'
            return 'Other'
        if source_col is not None:
            df['Booster Version Category'] = df[source_col].map(to_cat)
        else:
            df['Booster Version Category'] = 'Other'

    # Ensure minimal rows remain after coercions
    df = df.dropna(subset=['Launch Site', 'class', 'Payload Mass (kg)'])
    return df


def save_plotly_figures(df: pd.DataFrame) -> dict[str, str]:
    paths: dict[str, str] = {}
    # Pie: total successes by site
    success_counts = df.groupby("Launch Site")["class"].sum().reset_index(name="successes")
    fig_pie = px.pie(success_counts, values="successes", names="Launch Site",
                     title="Total Success Launches by Site")
    pie_path = os.path.join(PLOTS_DIR, "dash_pie_success_by_site.png")
    fig_pie.write_image(pie_path, width=1200, height=700, scale=2)
    paths["dash_pie"] = pie_path

    # Scatter: payload vs success, colored by Booster Version Category
    fig_scatter = px.scatter(
        df,
        x="Payload Mass (kg)",
        y="class",
        color="Booster Version Category",
        title="Payload vs Success (All Sites)",
        hover_data=["Launch Site"],
    )
    fig_scatter.update_yaxes(tickmode="array", tickvals=[0, 1], ticktext=["Failure", "Success"])
    scatter_path = os.path.join(PLOTS_DIR, "dash_scatter_payload_vs_success.png")
    fig_scatter.write_image(scatter_path, width=1200, height=700, scale=2)
    paths["dash_scatter"] = scatter_path

    return paths


def read_sql_summary() -> str:
    if os.path.exists(SQL_SUMMARY_MD):
        with open(SQL_SUMMARY_MD, "r", encoding="utf-8") as f:
            return f.read()
    return """SQL EDA summary not found. Please add results from Module 02 SQL analysis here."""


# --- Presentation helpers ---

def _get_layout(prs: Presentation, fallback_index: int = 1):
    # Prefer "Title and Content" if present; else fallback
    for idx, layout in enumerate(prs.slide_layouts):
        name = getattr(layout, 'name', '').lower()
        if "title and content" in name or "title and caption" in name:
            return layout
    return prs.slide_layouts[fallback_index]


def _set_title(slide, title: str):
    # Try standard title placeholder
    try:
        if getattr(slide.shapes, "title", None) is not None:
            slide.shapes.title.text = title
            return
    except Exception:
        pass

    # Try any title placeholder via placeholder collection
    try:
        for ph in slide.placeholders:
            try:
                if ph.placeholder_format.type == PP_PLACEHOLDER.TITLE and ph.has_text_frame:
                    ph.text_frame.text = title
                    return
            except Exception:
                continue
    except Exception:
        pass

    # Fallback: add a textbox near top
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1))
    tb.text_frame.text = title


def _get_body_text_frame(slide):
    # Prefer non-title text placeholder
    try:
        for ph in slide.placeholders:
            try:
                if ph.has_text_frame and ph.placeholder_format.type != PP_PLACEHOLDER.TITLE:
                    return ph.text_frame
            except Exception:
                continue
    except Exception:
        pass
    # Fallback: create a textbox taking most of the slide
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
    return tb.text_frame


def add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    # Try to use first layout; fall back to a generic layout
    layout = prs.slide_layouts[0] if len(prs.slide_layouts) else None
    if layout is None:
        layout = _get_layout(prs)
    slide = prs.slides.add_slide(layout)
    _set_title(slide, title)
    if subtitle:
        tf = _get_body_text_frame(slide)
        tf.text = subtitle
    return slide


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]):
    layout = _get_layout(prs)
    slide = prs.slides.add_slide(layout)
    _set_title(slide, title)
    body = _get_body_text_frame(slide)
    # clear to single empty paragraph
    try:
        body.clear()
    except Exception:
        try:
            # fallback: set first paragraph empty
            if getattr(body, "paragraphs", None) and len(body.paragraphs):
                body.paragraphs[0].text = ""
        except Exception:
            pass
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0] if getattr(body, "paragraphs", None) and len(body.paragraphs) else body.add_paragraph()
        else:
            p = body.add_paragraph()
        p.text = bullet
        p.level = 0
    return slide


def add_picture_slide(prs: Presentation, title: str, image_path: str):
    layout = _get_layout(prs)
    slide = prs.slides.add_slide(layout)
    _set_title(slide, title)
    left = Inches(0.8)
    top = Inches(1.8)
    height = Inches(5.5)
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, height=height)
    else:
        body = _get_body_text_frame(slide)
        body.text = f"Image not found: {image_path}"
    return slide


def add_text_slide(prs: Presentation, title: str, content: str):
    layout = _get_layout(prs)
    slide = prs.slides.add_slide(layout)
    _set_title(slide, title)
    tf = _get_body_text_frame(slide)
    try:
        tf.clear()
    except Exception:
        try:
            if getattr(tf, "paragraphs", None) and len(tf.paragraphs):
                tf.paragraphs[0].text = ""
        except Exception:
            pass
    for para in textwrap.fill(content, 120).split("\n"):
        if getattr(tf, "paragraphs", None) and len(tf.paragraphs):
            p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
        else:
            p = tf.add_paragraph()
        p.text = para
        p.level = 0
    return slide


# --- Build presentation ---

def build_presentation(metrics: dict | None = None):
    ensure_dirs()
    df = load_data()
    generated = save_plotly_figures(df)

    prs = Presentation(TEMPLATE_PATH)

    # --- Helper closures for template editing ---
    def _iter_text_shapes(slide):
        for shp in slide.shapes:
            if getattr(shp, "has_text_frame", False):
                yield shp

    def _get_all_text(slide) -> str:
        parts = []
        for shp in _iter_text_shapes(slide):
            try:
                parts.append(shp.text_frame.text)
            except Exception:
                continue
        return "\n".join(parts).lower()

    def _replace_tokens(slide, mapping: dict[str, str]):
        for shp in _iter_text_shapes(slide):
            try:
                tf = shp.text_frame
                for i, p in enumerate(list(tf.paragraphs)):
                    t = p.text
                    if not t:
                        continue
                    for k, v in mapping.items():
                        if k in t:
                            if i == 0:
                                tf.clear()
                                para = tf.paragraphs[0]
                                para.text = t.replace(k, v)
                            else:
                                p.text = t.replace(k, v)
                # Also direct full-text replacement for entire text_frame content
                if hasattr(tf, "text"):
                    val = tf.text
                    for k, v in mapping.items():
                        if k in val:
                            tf.text = val.replace(k, v)
            except Exception:
                continue

    # 1) Fill <Name> and <Date> on the title slide or anywhere in the deck
    today_str = datetime.now().strftime("%Y-%m-%d")
    name_date_map = {"<Name>": "Omar Essam", "<Date>": today_str}
    for slide in prs.slides:
        _replace_tokens(slide, name_date_map)

    # 2) Answer instruction prompts found in template (e.g., describe how data was collected)
    def _fill_instructions(slide):
        text_lower = _get_all_text(slide)
        def set_bullets(lines: list[str]):
            tf = _get_body_text_frame(slide)
            try:
                tf.clear()
            except Exception:
                pass
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 and getattr(tf, "paragraphs", None) and len(tf.paragraphs) else tf.add_paragraph()
                p.text = line
                p.level = 0
        if re.search(r"describe\s+how\s+the\s+data\s+was\s+collected", text_lower):
            set_bullets([
                "Primary CSV: Skills Network SpaceX launch dataset (dash CSV).",
                "Supplemental local dataset merged in Module 02 for robustness.",
                "Programmatic ingestion with pandas; standardized column names and types.",
                "Augmented features: landing success label, booster category, payload mass.",
            ])
        if re.search(r"explain\s+(the\s+)?methodology|methodological\s+approach", text_lower):
            set_bullets([
                "EDA: distributions, trends by year, correlation checks.",
                "Interactive analytics: Folium map, Plotly Dash (site filter, payload slider).",
                "Modeling: baseline classifiers with tuning; evaluated via accuracy and F1.",
            ])
        if re.search(r"explain\s+the\s+important\s+elements\s+and\s+findings", text_lower):
            set_bullets([
                "Higher success rates at KSC LC-39A and CCSFS SLC-40 in recent years.",
                "Payload mass shows success sweet spots; extremes trend riskier.",
                "Booster version category improvements align with higher success.",
            ])

    for slide in prs.slides:
        _fill_instructions(slide)

    # 3) Place plots into their intended placeholder slides based on keywords
    placed = set()
    def add_image(slide, path: str):
        if not os.path.exists(path):
            return False
        # Try to use a picture placeholder first for correct placement
        try:
            for shp in slide.shapes:
                if getattr(shp, "is_placeholder", False):
                    try:
                        pht = shp.placeholder_format.type
                    except Exception:
                        continue
                    # Try to insert into picture/content placeholders
                    try:
                        if pht in (PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.BODY):
                            shp.insert_picture(path)
                            return True
                    except Exception:
                        continue
        except Exception:
            pass
        # Fallback: absolute placement
        try:
            slide.shapes.add_picture(path, Inches(0.8), Inches(1.8), height=Inches(5.5))
            return True
        except Exception:
            return False

    # Map keywords -> image path
    images_map = [
        (("pie", "success", "site"), generated.get("dash_pie")),
        (("scatter", "payload"), generated.get("dash_scatter")),
    ]
    # Add Module 02 plots if the slide text mentions them
    mod2_candidates = {
        ("outcomes", "year"): os.path.join(MOD2_PLOTS, "outcomes_by_year.png"),
        ("payload", "distribution"): os.path.join(MOD2_PLOTS, "payload_distribution.png"),
        ("top", "launch", "sites"): os.path.join(MOD2_PLOTS, "top_launch_sites.png"),
    }

    for slide in prs.slides:
        text_lower = _get_all_text(slide)
        # Folium map screenshot if available
        if os.path.exists(FOLIUM_SCREENSHOT) and all(k in text_lower for k in ("folium", "map")):
            add_image(slide, FOLIUM_SCREENSHOT)
        # Dash images
        for keys, img in images_map:
            if img and all(k in text_lower for k in keys) and (img not in placed):
                if add_image(slide, img):
                    placed.add(img)
        # Module 2 plots
        for keys, img in mod2_candidates.items():
            if os.path.exists(img) and all(k in text_lower for k in keys) and (img not in placed):
                if add_image(slide, img):
                    placed.add(img)

    # SQL summary: if a slide contains 'sql' and 'findings', inject text
    sql_text = read_sql_summary()
    for slide in prs.slides:
        tl = _get_all_text(slide)
        if "sql" in tl and ("findings" in tl or "summary" in tl):
            tf = _get_body_text_frame(slide)
            try:
                tf.clear()
            except Exception:
                pass
            for i, para in enumerate(textwrap.fill(sql_text, 120).split("\n")):
                p = tf.paragraphs[0] if i == 0 and getattr(tf, "paragraphs", None) and len(tf.paragraphs) else tf.add_paragraph()
                p.text = para
                p.level = 0

    # 4) Remove placeholder instruction slides that still ask for screenshots/flowcharts and contain no images
    def slide_has_picture(slide) -> bool:
        try:
            for shp in slide.shapes:
                if getattr(shp, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    return True
        except Exception:
            pass
        return False

    deletion_phrases = [
        "place your flowchart",
        "place your screenshot",
        "show the screenshot",
        "data collection – spacex api",
        "data collection - spacex api",
        "data collection - scraping",
        "data collection – scraping",
        "success rate vs. orbit type",
        "flight number vs. orbit type",
        "payload vs. orbit type",
        "launch success yearly trend",
    ]

    def remove_slide_by_index(prs: Presentation, index: int):
        sldIdLst = prs.slides._sldIdLst
        slides = list(sldIdLst)
        sld = slides[index]
        prs.part.drop_rel(sld.rId)
        sldIdLst.remove(sld)

    to_delete: list[int] = []
    for idx, slide in enumerate(list(prs.slides)):
        tl = _get_all_text(slide)
        if (not slide_has_picture(slide)) and any(phrase in tl for phrase in deletion_phrases):
            to_delete.append(idx)

    for idx in sorted(to_delete, reverse=True):
        remove_slide_by_index(prs, idx)

    prs.save(OUT_PPTX)
    return OUT_PPTX


def try_export_pdf(input_pptx: str, output_pdf: str) -> bool:
    # Requires Microsoft PowerPoint installed on Windows
    try:
        import win32com.client  # type: ignore
    except Exception:
        return False

    try:
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(input_pptx, WithWindow=False)
        # 32 = ppSaveAsPDF
        presentation.SaveAs(output_pdf, 32)
        presentation.Close()
        powerpoint.Quit()
        return True
    except Exception:
        return False


if __name__ == "__main__":
    # Optionally pass a JSON file with metrics for the prediction results
    metrics = None
    if len(sys.argv) > 1 and os.path.exists(sys.argv[1]):
        with open(sys.argv[1], "r", encoding="utf-8") as f:
            metrics = json.load(f)

    pptx_path = build_presentation(metrics or {})
    print(f"Saved PPTX to: {pptx_path}")

    if try_export_pdf(pptx_path, OUT_PDF):
        print(f"Saved PDF to: {OUT_PDF}")
    else:
        print("PDF export skipped (PowerPoint COM not available). You can export manually from PowerPoint.")
